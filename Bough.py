
import os
import sys
import time
import json
import logging
import threading
import queue
import random
import re
import platform
import subprocess
import tempfile
import wave
import sqlite3
from datetime import datetime
from typing import Dict, List, Any, Optional

# Third-party imports with error handling
try:
    import pyttsx3
    import speech_recognition as sr
    import pyaudio
    import numpy as np
    import sympy as sp
    import requests
    from PIL import Image
    import pytesseract
    import PyPDF2
    HAS_DEPENDENCIES = True
except ImportError as e:
    print(f"Missing dependency: {e}")
    print("Please install required packages from requirements.txt")
    HAS_DEPENDENCIES = False

# Enhanced logging setup
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    handlers=[
        logging.FileHandler('bough_debug.log'),
        logging.StreamHandler(sys.stdout)
    ]
)
logger = logging.getLogger('JARVIS')

# Platform detection
is_mobile = platform.system() in ['Android', 'iOS'] or 'mobile' in platform.platform().lower()

# JARVIS Configuration
JARVIS_CONFIG = {
    'wake_word': 'jarvis',
    'voice_gender': 'male',
    'speech_rate': 180,
    'volume': 0.9,
    'response_style': 'professional',
    'lite_mode': False,
    'auto_start': False,
    'hotkey_enabled': True,
    'hotkey': 'ctrl+alt+j'
}

# Global state variables
wake_word_detected = False
listening_active = True
is_standby = False
race_mode_active = False
last_spoken_text = None
context_memory = []
command_history = []
knowledge_base = {}
alarm_time = None
todo_list = []
reminders = {}

def speak(text, voice_gender=None):
    global last_spoken_text

    # Avoid repeating the same text
    if text == last_spoken_text:
        return

    engine = pyttsx3.init()
    voices = engine.getProperty('voices')

    # Select voice based on gender preference or randomly
    selected_voice = None
    female_voices = [v for v in voices if "female" in v.name.lower()]
    male_voices = [v for v in voices if "male" in v.name.lower()]
    if voice_gender == "female":
        if female_voices:
            selected_voice = random.choice(female_voices)
    elif voice_gender == "male":
        if male_voices:
            selected_voice = random.choice(male_voices)
    else:
        # Randomly pick male or female voice if no preference
        all_voices = female_voices + male_voices
        if all_voices:
            selected_voice = random.choice(all_voices)

    if selected_voice:
        engine.setProperty('voice', selected_voice.id)
    else:
        engine.setProperty('voice', voices[0].id)

    # Optimize speech parameters for clarity and naturalness
    engine.setProperty('rate', 140)  # Slightly slower speech for clarity
    engine.setProperty('volume', 0.85)  # Slightly higher volume

    # JARVIS-style delivery with varied pauses
    sentences = [s.strip() for s in re.split(r'[.!?]', text) if s.strip()]
    for i, sentence in enumerate(sentences):
        print(f"VICK: {sentence}")
        engine.say(sentence)
        # Add varied pauses after certain statements
        if i < len(sentences)-1:
            if any(w in sentences[i].lower() for w in ['warning', 'alert', 'important']):
                engine.runAndWait()
                time.sleep(random.uniform(0.4, 0.7))
            else:
                engine.runAndWait()
                time.sleep(random.uniform(0.15, 0.35))
        else:
            engine.runAndWait()
            time.sleep(0.1)

    # Update the last spoken text
    last_spoken_text = text


import subprocess

def interpret_command(command):
    command = command.lower().strip()
    if any(exit_word in command for exit_word in ['exit', 'quit', 'stop', 'bye']):
        speak("Goodbye!")
        return False
    elif command == 'true':
        return 'command_true'
    elif command == 'false':
        return 'command_false'
    elif 'open notepad' in command:
        speak("Opening Notepad")
        subprocess.Popen(['notepad.exe'])
    elif 'open calculator' in command:
        speak("Opening Calculator")
        subprocess.Popen(['calc.exe'])
    elif 'what time' in command or 'current time' in command:
        now = datetime.now().strftime("%I:%M %p")
        speak(f"The current time is {now}")
    elif 'tell me a joke' in command or 'joke' in command:
        jokes = [
            "Why don't scientists trust atoms? Because they make up everything!",
            "Why did the computer show up at work late? It had a hard drive!",
            "Why do programmers prefer dark mode? Because light attracts bugs!"
        ]
        speak(random.choice(jokes))
    else:
        speak(f"You said: {command}")
    return True

def listen():
    """Listen for voice input and return recognized text using Whisper"""
    import whisper
    import tempfile
    import pyaudio
    import wave

    # Initialize Whisper model
    model = whisper.load_model("base")

    if is_mobile:
        try:
            text = input("Your command: ")
            return text.lower()
        except Exception as e:
            speak(f"Error reading input: {e}")
            return None

    try:
        CHUNK = 2048
        FORMAT = pyaudio.paInt16
        CHANNELS = 1
        RATE = 44100
        RECORD_SECONDS = 7

        p = pyaudio.PyAudio()
        stream = p.open(format=FORMAT,
                        channels=CHANNELS,
                        rate=RATE,
                        input=True,
                        frames_per_buffer=CHUNK)

        frames = []
        for _ in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
            data = stream.read(CHUNK)
            frames.append(data)

        stream.stop_stream()
        stream.close()
        p.terminate()

        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
            wf = wave.open(f.name, 'wb')
            wf.setnchannels(CHANNELS)
            wf.setsampwidth(p.get_sample_size(FORMAT))
            wf.setframerate(RATE)
            wf.writeframes(b''.join(frames))
            wf.close()

            audio_file = f.name

        # Use Whisper for recognition
        result = model.transcribe(audio_file)
        text = result["text"].strip().lower()
        print(f"You said: {text}")
        return text
    except Exception as e:
        speak(f"Error in listening: {e}")
        return None



def main():
    global is_standby
    speak("Hello! I am ready to assist you. Say 'exit' to quit.")
    running = True
    while running:
        try:
            command = listen()
        except Exception as e:
            speak(f"Error reading input: {e}")
            continue
        if command:
            action = interpret_command(command)

            if action == 'open':
                app_name = command.replace("open", "").strip()
                if app_name:
                    open_application(app_name)
            elif action == 'set_alarm':
                alarm_time_str = listen()
                if alarm_time_str:
                    set_alarm(alarm_time_str)
            elif action == 'add_task':
                task = listen()
                if task:
                    add_task(task)
            elif action == 'show_tasks':
                show_tasks()
            elif action == 'set_reminder':
                reminder_details = listen()
                if reminder_details:
                    set_reminder(reminder_details, "12:00 PM")
            elif action == 'tell_joke':
                tell_joke()
            elif action == 'get_time':
                speak(f"The time is {datetime.now().strftime('%I:%M %p')}")
            elif action == 'get_date':
                speak(f"Today is {datetime.now().strftime('%B %d, %Y')}")
            elif action == 'calculator':
                calculator()
            elif action == 'formula_one_updates':
                get_formula_one_updates()
            elif action == 'standby':
                is_standby = True
            elif action == 'wake_up':
                is_standby = False
            elif action == 'command_true':
                speak("Command acknowledged and executed.")
            elif action == 'command_false':
                speak("Command declined.")
            elif action == 'help':
                speak("I can help you with tasks, reminders, jokes, calculations, and more.")
            elif action == 'describe_personality':
                desc = personality_manager.describe_personality()
                speak(desc)
            elif action == 'open_audacity':
                open_audacity()
            elif action is None:
                fallback_response = respond(command if command else "")
                speak(fallback_response)
            else:
                speak(f"Command {action} is not implemented yet.")
        else:
            speak("I didn't catch that. Please repeat.")


if __name__ == "__main__":
    main()

import speech_recognition as sr

# Detect if running on mobile platform (simple heuristic)
is_mobile = platform.system() in ['Android', 'iOS'] or 'mobile' in platform.platform().lower()

# Model management toggling and lite mode flag
lite_mode = False

def detect_low_memory_system():
    """
    Detect if the system has low memory.
    This is a placeholder using platform and psutil if available.
    """
    try:
        import psutil
        mem = psutil.virtual_memory()
        # Consider low memory if total RAM less than 4GB
        if mem.total < 4 * 1024 * 1024 * 1024:
            return True
        return False
    except ImportError:
        # psutil not installed, fallback to platform check
        if is_mobile:
            return True
        return False

if not is_mobile:
    import pyttsx3
    import cv2
    import subprocess
    import sounddevice as sd
else:
    # Mobile stubs or alternative imports
    pyttsx3 = None
    cv2 = None
    subprocess = None
    sounddevice = None

# Stub classes for new advanced features

class SystemMonitor:
    def __init__(self):
        self.monitoring = False

    def start_monitoring(self):
        self.monitoring = True
        # Placeholder: Start system health monitoring in a separate thread
        threading.Thread(target=self._monitor_loop, daemon=True).start()

    def _monitor_loop(self):
        while self.monitoring:
            # Placeholder: Check system health and perform self-healing if needed
            time.sleep(10)

    def stop_monitoring(self):
        self.monitoring = False

class PredictiveIntelligence:
    def __init__(self):
        pass

    def predict_next(self, command, context):
        # Placeholder: Return a prediction string or None
        return None

class WearableIntegration:
    def __init__(self):
        pass

    def get_wearable_data(self):
        # Placeholder: Return data from connected wearables
        return {}

    def send_command(self, command):
        # Placeholder: Send command to wearable device
        pass

import tempfile

class DeveloperAssistant:
    def __init__(self):
        pass

    def generate_code(self, prompt, language=None):
        """
        Generate code based on prompt and optional language.
        This is a stub method; integration with AI model can be added.
        """
        # For now, return a placeholder string
        return f"# Code generated for prompt: {prompt}\n# Language: {language or 'unspecified'}\n"

    def run_code(self, code, language):
        """
        Run code snippet in the specified language.
        Supports common languages: python, javascript, java, c, cpp, bash.
        Returns output and error messages.
        """
        import subprocess
        import os

        with tempfile.TemporaryDirectory() as tmpdir:
            file_map = {
                'python': 'script.py',
                'javascript': 'script.js',
                'js': 'script.js',
                'java': 'Main.java',
                'c': 'program.c',
                'cpp': 'program.cpp',
                'bash': 'script.sh',
                'sh': 'script.sh'
            }
            lang = language.lower()
            if lang not in file_map:
                return "", f"Language '{language}' not supported for execution."

            file_name = file_map[lang]
            file_path = os.path.join(tmpdir, file_name)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(code)

            try:
                if lang == 'python':
                    cmd = ['python', file_path]
                elif lang in ['javascript', 'js']:
                    cmd = ['node', file_path]
                elif lang == 'java':
                    # Compile and run Java
                    compile_proc = subprocess.run(['javac', file_path], capture_output=True, text=True)
                    if compile_proc.returncode != 0:
                        return "", f"Compilation error:\n{compile_proc.stderr}"
                    cmd = ['java', '-cp', tmpdir, 'Main']
                elif lang == 'c':
                    exe_path = os.path.join(tmpdir, 'program.exe')
                    compile_proc = subprocess.run(['gcc', file_path, '-o', exe_path], capture_output=True, text=True)
                    if compile_proc.returncode != 0:
                        return "", f"Compilation error:\n{compile_proc.stderr}"
                    cmd = [exe_path]
                elif lang == 'cpp':
                    exe_path = os.path.join(tmpdir, 'program.exe')
                    compile_proc = subprocess.run(['g++', file_path, '-o', exe_path], capture_output=True, text=True)
                    if compile_proc.returncode != 0:
                        return "", f"Compilation error:\n{compile_proc.stderr}"
                    cmd = [exe_path]
                elif lang in ['bash', 'sh']:
                    cmd = ['bash', file_path]
                else:
                    return "", f"Execution for language '{language}' not implemented."

                proc = subprocess.run(cmd, capture_output=True, text=True, timeout=30)
                return proc.stdout, proc.stderr
            except subprocess.TimeoutExpired:
                return "", "Execution timed out."
            except Exception as e:
                return "", f"Error during execution: {e}"

    def debug_code(self, code, language):
        """
        Debug code by checking syntax errors or running linters.
        Returns diagnostics or error messages.
        """
        import subprocess
        import os

        lang = language.lower()
        with tempfile.TemporaryDirectory() as tmpdir:
            file_map = {
                'python': 'script.py',
                'javascript': 'script.js',
                'js': 'script.js',
                'java': 'Main.java',
                'c': 'program.c',
                'cpp': 'program.cpp',
                'bash': 'script.sh',
                'sh': 'script.sh'
            }
            if lang not in file_map:
                return f"Language '{language}' not supported for debugging."

            file_name = file_map[lang]
            file_path = os.path.join(tmpdir, file_name)
            with open(file_path, 'w', encoding='utf-8') as f:
                f.write(code)

            try:
                if lang == 'python':
                    # Use python -m py_compile for syntax check
                    proc = subprocess.run(['python', '-m', 'py_compile', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No syntax errors detected."
                    else:
                        return proc.stderr
                elif lang in ['javascript', 'js']:
                    # Use eslint if available
                    proc = subprocess.run(['eslint', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No linting errors detected."
                    else:
                        return proc.stdout + proc.stderr
                elif lang == 'java':
                    proc = subprocess.run(['javac', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No compilation errors detected."
                    else:
                        return proc.stderr
                elif lang == 'c':
                    proc = subprocess.run(['gcc', '-fsyntax-only', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No syntax errors detected."
                    else:
                        return proc.stderr
                elif lang == 'cpp':
                    proc = subprocess.run(['g++', '-fsyntax-only', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No syntax errors detected."
                    else:
                        return proc.stderr
                elif lang in ['bash', 'sh']:
                    proc = subprocess.run(['bash', '-n', file_path], capture_output=True, text=True)
                    if proc.returncode == 0:
                        return "No syntax errors detected."
                    else:
                        return proc.stderr
                else:
                    return f"Debugging for language '{language}' not implemented."
            except Exception as e:
                return f"Error during debugging: {e}"

    def check_errors(self, code, language):
        """
        Analyze code for errors by attempting to run debug_code.
        Returns error messages or empty string if no errors.
        """
        diagnostics = self.debug_code(code, language)
        if "No" in diagnostics:
            return ""
        else:
            return diagnostics

class DataFusionManager:
    def __init__(self):
        pass

    def fuse_data(self, data_sources):
        # Placeholder: Combine data from multiple sources
        return {}

class ModularExpansionManager:
    def __init__(self):
        self.modules = []

    def load_module(self, module):
        # Placeholder: Load a modular expansion
        self.modules.append(module)

    def unload_module(self, module):
        # Placeholder: Unload a modular expansion
        if module in self.modules:
            self.modules.remove(module)

class VoiceSynthesisCustomizer:
    def __init__(self):
        pass

    def customize_voice(self, settings):
        # Placeholder: Customize voice synthesis parameters
        pass

class CollaborationManager:
    def __init__(self):
        pass

    def start_session(self, users):
        # Placeholder: Start collaborative session
        pass

    def sync_data(self, data):
        # Placeholder: Sync data among collaborators
        pass

class OfflineLocationDetector:
    def __init__(self):
        pass

    def get_location(self):
        # Placeholder: Return location data without internet
        return None


class OfflineLocationDetector:
    def __init__(self):
        pass

    def get_location(self):
        # Placeholder: Return location data without internet
        return None



# Cloud Sync Configuration
CLOUD_SYNC_ENABLED = False
CLOUD_CONFIG = {
    'api_key': '',
    'cloud_url': 'https://api.example.com'
}

# --- Sim Racing Assistant class added ---
class SimRacingAssistant:
    def __init__(self):
        # Example car setups for different tracks
        self.car_setups = {
            "monza": {"suspension": "soft", "aero": "low", "brakes": "medium"},
            "spa": {"suspension": "medium", "aero": "medium", "brakes": "high"},
            "silverstone": {"suspension": "hard", "aero": "high", "brakes": "medium"},
        }
        # Tyre types and their best use conditions
        self.tyre_types = {
            "soft": "best for qualifying and short stints",
            "medium": "balanced for race distance",
            "hard": "best for long stints and durability",
            "wet": "use in wet conditions",
            "intermediate": "use in light rain or drying track"
        }
        # Preset strategy profiles
        self.strategy_profiles = {
            "aggressive": "Push hard, prioritize overtaking, and minimize pit stops.",
            "conservative": "Focus on tire preservation and consistent lap times.",
            "wet_specialist": "Adapt quickly to changing wet conditions and use wet tires effectively."
        }

    def get_car_setup(self, track_name):
        track = track_name.lower()
        setup = self.car_setups.get(track)
        if setup:
            return f"For {track_name.title()}, use suspension: {setup['suspension']}, aero: {setup['aero']}, brakes: {setup['brakes']}."
        else:
            return f"No specific car setup found for {track_name.title()}. Use a balanced setup."

    def get_tyre_recommendation(self, weather_condition):
        condition = weather_condition.lower()
        if "wet" in condition or "rain" in condition:
            return "Use wet tyres if the track is wet, or intermediate tyres if the track is drying."
        elif "dry" in condition or "clear" in condition:
            return "Use soft, medium, or hard tyres depending on race length and strategy."
        else:
            return "Use medium tyres as a safe default."

    def get_gap_info(self, gap_front, gap_behind):
        return f"The gap to the car in front is {gap_front} seconds, and the gap to the car behind is {gap_behind} seconds."

    def get_race_strategy(self, laps_remaining, current_tyre, fuel_level):
        strategy = "Maintain consistent lap times."
        if fuel_level < 10:
            strategy += " Consider a pit stop soon for refueling."
        if current_tyre.lower() == "soft" and laps_remaining > 10:
            strategy += " You might want to switch to harder tyres for durability."
        if laps_remaining <= 5:
            strategy += " Push hard to gain positions in the final laps."
        return strategy

    def race_engineer_mode(self, telemetry_data):
        """
        Analyze telemetry data for car performance, tire degradation, fuel usage, and pit stop advice.
        :param telemetry_data: dict with keys like 'tire_wear', 'fuel_level', 'lap_times'
        :return: string advice
        """
        tire_wear = telemetry_data.get('tire_wear', 0)
        fuel_level = telemetry_data.get('fuel_level', 100)
        lap_times = telemetry_data.get('lap_times', [])
        advice = []

        if tire_wear > 70:
            advice.append("Tire wear is high, consider a pit stop soon.")
        else:
            advice.append("Tire wear is within acceptable limits.")

        if fuel_level < 15:
            advice.append("Fuel level is low, plan for refueling.")
        else:
            advice.append("Fuel level is sufficient for now.")

        if lap_times:
            avg_lap = sum(lap_times) / len(lap_times)
            advice.append(f"Average lap time is {avg_lap:.2f} seconds.")

        return " ".join(advice)

    def strategy_advisor(self, track_conditions, competitor_tactics, driver_preferences):
        """
        Suggest race strategies based on track conditions, competitor tactics, and driver preferences.
        :return: string strategy advice
        """
        weather = track_conditions.get('weather', 'dry')
        strategy = "Base strategy: Maintain consistent pace."

        if weather == 'wet':
            strategy += " Use wet tyres and be cautious in corners."
        else:
            strategy += " Use optimal tyre strategy based on race length."

        if competitor_tactics == 'aggressive':
            strategy += " Be prepared for defensive driving."

        if driver_preferences == 'conservative':
            strategy += " Focus on tire preservation."

        return strategy

    def team_management_assistant(self, career_data):
        """
        Assist with contract negotiations, driver selection, car upgrades, and budget management.
        :param career_data: dict with team info
        :return: string advice
        """
        budget = career_data.get('budget', 1000000)
        upgrades = career_data.get('upgrades', [])
        advice = []

        if budget < 100000:
            advice.append("Budget is tight, prioritize essential upgrades.")
        else:
            advice.append("Consider investing in aerodynamic upgrades.")

        if 'engine' not in upgrades:
            advice.append("Engine upgrade is recommended.")

        return " ".join(advice)

    def telemetry_analyzer(self, telemetry_data):
        """
        Analyze lap times, tire temps, sector speeds, and braking efficiency.
        :param telemetry_data: dict with telemetry info
        :return: string analysis
        """
        lap_times = telemetry_data.get('lap_times', [])
        tire_temps = telemetry_data.get('tire_temps', [])
        sector_speeds = telemetry_data.get('sector_speeds', [])
        braking_efficiency = telemetry_data.get('braking_efficiency', 0)

        analysis = []

        if lap_times:
            best_lap = min(lap_times)
            analysis.append(f"Best lap time: {best_lap:.2f} seconds.")

        if tire_temps:
            avg_temp = sum(tire_temps) / len(tire_temps)
            analysis.append(f"Average tire temperature: {avg_temp:.1f}°C.")

        if sector_speeds:
            avg_speed = sum(sector_speeds) / len(sector_speeds)
            analysis.append(f"Average sector speed: {avg_speed:.1f} km/h.")

        analysis.append(f"Braking efficiency: {braking_efficiency}%.")

        return " ".join(analysis)

    def live_radio_communication(self, message_type):
        """
        Simulate interactive team radio messages.
        :param message_type: string indicating message context
        :return: string message
        """
        messages = {
            "push": "Push hard now, you have a gap to the car behind.",
            "conserve": "Conserve your tires and fuel for the next stint.",
            "pit_soon": "Prepare for a pit stop on the next lap.",
            "weather_alert": "Weather is changing, adjust your strategy accordingly."
        }
        return messages.get(message_type, "Keep up the good work, stay focused.")

    def dynamic_weather_adaptation(self, weather_conditions):
        """
        Offer adaptive strategies based on changing track conditions.
        :param weather_conditions: dict with weather info
        :return: string advice
        """
        weather = weather_conditions.get('weather', 'dry')
        temperature = weather_conditions.get('temperature', 25)
        advice = "Monitor weather conditions closely."

        if weather == 'rain':
            advice = "Switch to wet tyres and drive cautiously."
        elif weather == 'dry' and temperature > 30:
            advice = "Consider harder tyres to manage heat."

        return advice

    def opponent_prediction_system(self, past_race_data):
        """
        Predict competitor behavior using past race data.
        :param past_race_data: list of dicts with competitor info
        :return: string prediction
        """
        # Placeholder simple prediction
        if not past_race_data:
            return "No past race data available for predictions."

        aggressive_drivers = [d for d in past_race_data if d.get('aggressiveness', 0) > 7]
        if aggressive_drivers:
            return f"Expect aggressive moves from {len(aggressive_drivers)} competitors."
        else:
            return "Competitors are expected to race conservatively."

    def pit_wall_decision_center(self, race_status):
        """
        Simulate pit wall decisions like undercut attempts, double-stack strategies, and fuel-saving.
        :param race_status: dict with race info
        :return: string decision advice
        """
        laps_to_go = race_status.get('laps_to_go', 20)
        pit_stop_window = race_status.get('pit_stop_window', (10, 15))
        fuel_level = race_status.get('fuel_level', 50)
        advice = "Maintain current strategy."

        if laps_to_go in range(pit_stop_window[0], pit_stop_window[1]+1):
            advice = "Consider an undercut pit stop to gain track position."
        if fuel_level < 15:
            advice += " Save fuel to avoid extra pit stops."

        return advice

    def assist(self, query, **kwargs):
        query = query.lower()
        if "car setup" in query:
            track = kwargs.get("track", "unknown track")
            return self.get_car_setup(track)
        elif "tyre" in query or "tire" in query:
            weather = kwargs.get("weather", "dry")
            return self.get_tyre_recommendation(weather)
        elif "gap" in query:
            gap_front = kwargs.get("gap_front", "unknown")
            gap_behind = kwargs.get("gap_behind", "unknown")
            return self.get_gap_info(gap_front, gap_behind)
        elif "strategy" in query:
            laps_remaining = kwargs.get("laps_remaining", 20)
            current_tyre = kwargs.get("current_tyre", "medium")
            fuel_level = kwargs.get("fuel_level", 50)
            return self.get_race_strategy(laps_remaining, current_tyre, fuel_level)
        elif "engineer" in query:
            telemetry_data = kwargs.get("telemetry_data", {})
            return self.race_engineer_mode(telemetry_data)
        elif "advisor" in query:
            track_conditions = kwargs.get("track_conditions", {})
            competitor_tactics = kwargs.get("competitor_tactics", "normal")
            driver_preferences = kwargs.get("driver_preferences", "balanced")
            return self.strategy_advisor(track_conditions, competitor_tactics, driver_preferences)
        elif "team" in query:
            career_data = kwargs.get("career_data", {})
            return self.team_management_assistant(career_data)
        elif "telemetry" in query:
            telemetry_data = kwargs.get("telemetry_data", {})
            return self.telemetry_analyzer(telemetry_data)
        elif "radio" in query:
            message_type = kwargs.get("message_type", "default")
            return self.live_radio_communication(message_type)
        elif "weather" in query:
            weather_conditions = kwargs.get("weather_conditions", {})
            return self.dynamic_weather_adaptation(weather_conditions)
        elif "opponent" in query:
            past_race_data = kwargs.get("past_race_data", [])
            return self.opponent_prediction_system(past_race_data)
        elif "pit wall" in query:
            race_status = kwargs.get("race_status", {})
            return self.pit_wall_decision_center(race_status)
        elif "preset strategy" in query:
            profile = kwargs.get("profile", "balanced")
            return self.strategy_profiles.get(profile, "No such strategy profile found.")
        else:
            return "I can assist with car setup, tyre recommendations, gap info, race strategy, engineering, team management, telemetry analysis, live radio, weather adaptation, opponent prediction, and pit wall decisions. Please specify your request."

# Smart Home Integration
SMART_HOME_DEVICES = {
    'lights': {
        'bedroom': '192.168.1.100',
        'living_room': '192.168.1.101'
    },
    'thermostat': '192.168.1.102'
}

# Initialize enhanced components

# Enhanced audio settings
SAMPLE_RATE = 16000
CHANNELS = 1
BLOCKSIZE = 8192
DEVICE = None  # Default audio device

import numpy as np
import wave
import pyaudio
import tempfile
import pyttsx3
import tkinter as tk
import random
import threading
import queue

# GUI class for Bough AI Assistant
class BoughGUI:
    def __init__(self, root, command_queue, response_queue):
        self.root = root
        self.command_queue = command_queue
        self.response_queue = response_queue
        self.root.title("Bough AI Assistant")
        self.root.geometry("600x400")
        self.canvas = tk.Canvas(root, width=200, height=200)
        self.canvas.pack(side=tk.LEFT, padx=10, pady=10)
        self.current_form = "robot"
        self.expression = "neutral"
        self.forms = {
            "robot": self.draw_robot,
            "human": self.draw_human,
            "animal": self.draw_animal
        }
        self.draw_face()
        self.blink()
        self.move()
        self.text_area = tk.Text(root, wrap=tk.WORD, state=tk.DISABLED)
        self.text_area.pack(padx=10, pady=10, fill=tk.BOTH, expand=True)

        self.entry = tk.Entry(root)
        self.entry.pack(padx=10, pady=(0,10), fill=tk.X)
        self.entry.bind("<Return>", self.send_command)

        self.send_button = tk.Button(root, text="Send", command=self.send_command)
        self.send_button.pack(padx=10, pady=(0,10))

        self.update_text_area()

    def send_command(self, event=None):
        command = self.entry.get().strip()
        if command:
            self.append_text(f"You: {command}\n")
            self.command_queue.put(command)
            self.entry.delete(0, tk.END)

    def append_text(self, text):
        self.text_area.config(state=tk.NORMAL)
        self.text_area.insert(tk.END, text)
        self.text_area.see(tk.END)
        self.text_area.config(state=tk.DISABLED)

    def update_text_area(self):
        try:
            while True:
                response = self.response_queue.get_nowait()
                self.append_text(f"Bough: {response}\n")
        except queue.Empty:
            pass
        self.root.after(100, self.update_text_area)


def listen():
    import speech_recognition as sr
    r = sr.Recognizer()
    if is_mobile:
        # Removed speaking back for typing input
        try:
            text = input("Your command: ")
            return text.lower()
        except Exception as e:
            # Removed speaking back for error
            return None

    try:
        CHUNK = 2048
        FORMAT = pyaudio.paInt16
        CHANNELS = 1
        RATE = 44100
        RECORD_SECONDS = 7

        p = pyaudio.PyAudio()
        stream = p.open(format=FORMAT,
                        channels=CHANNELS,
                        rate=RATE,
                        input=True,
                        frames_per_buffer=CHUNK)

        frames = []
        for _ in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
            data = stream.read(CHUNK)
            frames.append(data)

        stream.stop_stream()
        stream.close()
        p.terminate()

        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
            wf = wave.open(f.name, 'wb')
            wf.setnchannels(CHANNELS)
            wf.setsampwidth(p.get_sample_size(FORMAT))
            wf.setframerate(RATE)
            wf.writeframes(b''.join(frames))
            wf.close()

            audio_file = f.name

        with sr.AudioFile(audio_file) as source:
            audio = r.record(source)
        try:
            text = r.recognize_google(audio)
            # Removed speaking back of recognized text
            return text.lower()
        except sr.UnknownValueError:
            # Removed speaking back for error
            return None
        except sr.RequestError as e:
            # Removed speaking back for error
            return None
        except Exception as e:
            # Removed speaking back for error
            return None
    except Exception as e:
        # Removed speaking back for error
        return None

# Initialize TTS with better voice
try:
    engine = pyttsx3.init()
    voices = engine.getProperty('voices')
    # Choose a faster voice if available
    fast_voice = None
    for voice in voices:
        if "female" in voice.name.lower() or "default" in voice.name.lower():
            fast_voice = voice.id
            break
    if fast_voice:
        engine.setProperty('voice', fast_voice)
    else:
        engine.setProperty('voice', voices[0].id)
    engine.setProperty('rate', 200)  # Faster speech rate
except Exception as e:
    print(f"Error initializing TTS engine: {e}")
    engine = None


alarm_time = None
todo_list = []
reminders = {}
is_standby = False
last_spoken_text = None
context_memory = []
command_history = []
knowledge_base = {}
knowledge_file = "knowledge_base.json"

# Persistent local memory using SQLite
import sqlite3

def init_local_memory_db():
    conn = sqlite3.connect('bough_local_memory.db')
    c = conn.cursor()
    c.execute('''
        CREATE TABLE IF NOT EXISTS user_preferences (
            key TEXT PRIMARY KEY,
            value TEXT
        )
    ''')
    c.execute('''
        CREATE TABLE IF NOT EXISTS voice_memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            voice_data BLOB,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    c.execute('''
        CREATE TABLE IF NOT EXISTS mood_memory (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            mood TEXT,
            timestamp DATETIME DEFAULT CURRENT_TIMESTAMP
        )
    ''')
    conn.commit()
    conn.close()

def save_user_preference(key, value):
    conn = sqlite3.connect('bough_local_memory.db')
    c = conn.cursor()
    c.execute('REPLACE INTO user_preferences (key, value) VALUES (?, ?)', (key, value))
    conn.commit()
    conn.close()

def get_user_preference(key):
    conn = sqlite3.connect('bough_local_memory.db')
    c = conn.cursor()
    c.execute('SELECT value FROM user_preferences WHERE key = ?', (key,))
    row = c.fetchone()
    conn.close()
    return row[0] if row else None

# Initialize local memory DB on startup
init_local_memory_db()

# Load knowledge base from file
try:
    with open(knowledge_file, 'r') as f:
        knowledge_base = json.load(f)
except (FileNotFoundError, json.JSONDecodeError):
    knowledge_base = {}

import socket

def is_connected(host="8.8.8.8", port=53, timeout=3):
    """
    Check internet connectivity by attempting to connect to a DNS server.
    """
    try:
        socket.setdefaulttimeout(timeout)
        socket.socket(socket.AF_INET, socket.SOCK_STREAM).connect((host, port))
        return True
    except Exception:
        return False

def speak(text, priority="normal", voice_gender=None):
    global last_spoken_text

    # Avoid repeating the same text
    if text == last_spoken_text and priority != "alert":
        return

    if engine is None:
        # Mobile or no TTS engine: fallback to print
        print(f"Speak: {text}")
        last_spoken_text = text
        return

    try:


        # Select voice based on gender preference or randomly
        voices = engine.getProperty('voices')
        selected_voice = None

        if voice_gender == "female":
            female_voices = [v for v in voices if "female" in v.name.lower()]
            if female_voices:
                selected_voice = random.choice(female_voices)
        elif voice_gender == "male":
            male_voices = [v for v in voices if "male" in v.name.lower()]
            if male_voices:
                # Prefer JARVIS-style male voice if available
                jarvis_voices = [v for v in male_voices if "jarvis" in v.name.lower()]
                if jarvis_voices:
                    selected_voice = random.choice(jarvis_voices)
                else:
                    selected_voice = random.choice(male_voices)
        else:
            # Randomly pick male or female voice if no preference
            male_voices = [v for v in voices if "male" in v.name.lower()]
            female_voices = [v for v in voices if "female" in v.name.lower()]
            all_voices = male_voices + female_voices
            if all_voices:
                selected_voice = random.choice(all_voices)

        if selected_voice:
            engine.setProperty('voice', selected_voice.id)
        else:
            # Fallback to first voice
            engine.setProperty('voice', voices[0].id)

        # Optimize speech parameters for clarity and naturalness
        engine.setProperty('rate', 140)  # Slightly slower speech for clarity
        engine.setProperty('volume', 0.85)  # Slightly higher volume

        # JARVIS-style delivery with varied pauses
        sentences = [s.strip() for s in re.split(r'[.!?]', text) if s.strip()]
        for i, sentence in enumerate(sentences):
            print(f"VICK: {sentence}")
            engine.say(sentence)
            # Add varied pauses after certain statements
            if i < len(sentences)-1:
                if any(w in sentences[i].lower() for w in ['warning', 'alert', 'important']):
                    engine.runAndWait()
                    time.sleep(random.uniform(0.4, 0.7))
                else:
                    engine.runAndWait()
                    time.sleep(random.uniform(0.15, 0.35))
            else:
                engine.runAndWait()
                time.sleep(0.1)

        # Update the last spoken text
        last_spoken_text = text

    except Exception as e:
        print(f"Error in speech synthesis: {str(e)}")

import whisper
import speech_recognition as sr
import tempfile
import pyaudio
import wave
import logging
import os
import string
import glob
import json
import platform
import socket
import subprocess
import time
import random
import re
import threading
import queue
import hashlib
import PyPDF2
import requests
import numpy as np


# Define missing variables and functions

# Define drives for search (e.g., Windows drives)
drives = ['C:\\', 'D:\\', 'E:\\']  # Adjust as per system

# Define exclude_dirs for search
exclude_dirs = ['Windows', 'Program Files', 'Program Files (x86)', 'AppData', 'Temp']

# Implement search_offline_content as a placeholder
def search_offline_content(query):
    # Placeholder implementation: return no results
    return "Search functionality is not implemented yet."

# Initialize logger
logger = logging.getLogger('BoughAI')
logger.setLevel(logging.INFO)
handler = logging.StreamHandler()
formatter = logging.Formatter('%(asctime)s - %(name)s - %(levelname)s - %(message)s')
handler.setFormatter(formatter)
logger.addHandler(handler)

# Lazy load Whisper model
_whisper_model = None

def transcribe_with_whisper(audio_path):
    global _whisper_model
    if _whisper_model is None:
        _whisper_model = whisper.load_model("base")
    result = _whisper_model.transcribe(audio_path, fp16=False)
    return result["text"].strip().lower()

import platform

# Enhanced listen function with multiple voice recognition options and enhancements

import numpy as np
import wave


# Lazy load WhisperX and DeepSpeech models
whisperx_model = None
ds_model = None



def recognize_sphinx():
    r = sr.Recognizer()
    with sr.Microphone() as source:
        speak("Listening with Sphinx...")
        r.adjust_for_ambient_noise(source)
        audio = r.listen(source)
    try:
        text = r.recognize_sphinx(audio)
        speak(f"Sphinx recognized: {text}")
        return text.lower()
    except sr.UnknownValueError:
        speak("Sphinx could not understand audio.")
        return None
    except sr.RequestError as e:
        speak(f"Sphinx error; {e}")
        return None

def recognize_google_web(audio):
    r = sr.Recognizer()
    try:
        text = r.recognize_google(audio)
        speak(f"Google Web Speech recognized: {text}")
        return text.lower()
    except sr.UnknownValueError:
        speak("Google Web Speech could not understand audio.")
        return None
    except sr.RequestError as e:
        speak(f"Could not request results from Google Web Speech service; {e}")
        return None
    except Exception as e:
        speak(f"An unexpected error occurred in Google Web Speech: {e}")
        return None

import win32com.client

def listen_windows_sapi():
    try:
        recognizer = win32com.client.Dispatch("SAPI.SpInprocRecognizer")
        context = recognizer.CreateRecoContext()
        grammar = context.CreateGrammar()
        grammar.DictationSetState(1)  # Enable dictation
        result = None

        def callback(event):
            nonlocal result
            if event.Result:
                result = event.Result.PhraseInfo.GetText()

        context.EventInterests = 1  # SPEI_RECOGNITION
        context.Recognition += callback

        # Wait for recognition event or timeout
        import pythoncom
        import time
        timeout = 10  # seconds
        start_time = time.time()
        while result is None and (time.time() - start_time) < timeout:
            pythoncom.PumpWaitingMessages()
            time.sleep(0.1)

        grammar.DictationSetState(0)  # Disable dictation

        if result:
            speak(f"Windows SAPI recognized: {result}")
            return result.lower()
        else:
            speak("Windows SAPI did not recognize any speech.")
            return None
    except Exception as e:
        speak(f"Windows SAPI error: {e}")
        return None

# Function to open Audacity application
def open_audacity():
    """Open Audacity application if installed"""
    try:
        # Common install paths for Audacity on Windows
        possible_paths = [
            os.path.expandvars(r"%ProgramFiles%\Audacity\audacity.exe"),
            os.path.expandvars(r"%ProgramFiles(x86)%\Audacity\audacity.exe"),
            os.path.expandvars(r"%LOCALAPPDATA%\Programs\Audacity\audacity.exe")
        ]
        for path in possible_paths:
            if os.path.exists(path):
                subprocess.Popen([path])
                speak("Opening Audacity")
                return True
        # Try to open by command if in PATH
        subprocess.Popen(["audacity"])
        speak("Opening Audacity")
        return True
    except Exception as e:
        print(f"Error opening Audacity: {e}")
        speak("Could not open Audacity")
        return False

# Add 'open audacity' command to command_map in interpret_command

def respond(text):
    """Generate natural language responses"""
    responses = {
        "greeting": ["Hello there!", "Hi! How can I help?", "Greetings!"],
        "time": [f"The time is {datetime.now().strftime('%I:%M %p')}"],
        "date": [f"Today is {datetime.now().strftime('%B %d, %Y')}"],
        "default": ["I didn't quite catch that", "Could you repeat that?"]
    }
    
    # Simple intent detection
    if any(word in text for word in ['hi', 'hello', 'hey']):
        return random.choice(responses["greeting"])
    elif 'time' in text:
        return random.choice(responses["time"])
    elif 'date' in text:
        return random.choice(responses["date"])
    else:
        return random.choice(responses["default"])



# Function to open applications
def open_application(app_name):
    """Open any application by searching common install locations"""
    try:
        # Search common Windows program locations
        search_paths = [
            os.path.expandvars("%ProgramFiles%"),
            os.path.expandvars("%ProgramFiles(x86)%"),
            os.path.expandvars("%LOCALAPPDATA%\\Programs"),
            os.path.expandvars("%APPDATA%"),
            os.path.expandvars("%SystemRoot%\\System32")
        ]
        
        # Search for executable with matching name
        for root in search_paths:
            for dirpath, dirnames, filenames in os.walk(root):
                for filename in filenames:
                    if filename.lower().startswith(app_name.lower()):
                        ext = os.path.splitext(filename)[1].lower()
                        if ext in ['.exe', '.bat', '.cmd', '.msi']:
                            full_path = os.path.join(dirpath, filename)
                            subprocess.Popen([full_path])
                            speak(f"Opening {filename}")
                            return True
        
        # If not found, try direct execution (may work for system apps)
        try:
            subprocess.Popen([app_name])
            speak(f"Attempting to open {app_name}")
            return True
        except:
            speak(f"Could not find application: {app_name}")
            return False
            
    except Exception as e:
        print(f"Error opening application: {e}")
        speak("Sorry, I encountered an error while trying to open that application")
        return False

# Function to set an alarm
def set_alarm(alarm_time_str):
    global alarm_time
    try:
        alarm_time = datetime.strptime(alarm_time_str, "%I:%M %p")
        speak(f"Alarm set for {alarm_time.strftime('%I:%M %p')}.")
        threading.Thread(target=monitor_alarm).start()
    except ValueError:
        speak("Sorry, I didn't understand the time. Please say something like 'set alarm for 7:30 AM'.")

# Function to monitor the alarm
def monitor_alarm():
    global alarm_time
    while alarm_time:
        if datetime.now().strftime("%I:%M %p") == alarm_time.strftime("%I:%M %p"):
            speak("Wake up! Your alarm is ringing.")
            alarm_time = None
        time.sleep(10)

# Function to manage tasks
def add_task(task):
    if task:
        todo_list.append(task)
        speak(f"Task added: {task}")
    else:
        speak("Please enter a task name.")

def show_tasks():
    if todo_list:
        speak("Here are your tasks:")
        for i, task in enumerate(todo_list, 1):
            speak(f"Task {i}: {task}")
    else:
        speak("Your to-do list is empty.")

# Function to set reminders
def set_reminder(task, reminder_time_str):
    try:
        reminder_time = datetime.strptime(reminder_time_str, "%I:%M %p")
        reminders[task] = reminder_time
        speak(f"Reminder set for {task} at {reminder_time.strftime('%I:%M %p')}.")
        threading.Thread(target=monitor_reminder, args=(task, reminder_time)).start()
    except ValueError:
        speak("Sorry, I didn't understand the time. Please say something like 'remind me to call John at 3:00 PM'.")

# Function to monitor reminders
def monitor_reminder(task, reminder_time):
    while task in reminders:
        if datetime.now().strftime("%I:%M %p") == reminder_time.strftime("%I:%M %p"):
            speak(f"Reminder: {task}")
            del reminders[task]
        time.sleep(10)

# Function to send email
def send_email(to, subject, body):
    speak("Email functionality is not available in offline mode.")

# Function to add calendar event
def add_calendar_event(summary, start_time, end_time):
    speak("Calendar functionality is not available in offline mode.")

# Function to tell a joke
def tell_joke():
    jokes = [
        "Why don't scientists trust atoms? Because they make up everything!",
        "Did you hear about the mathematician who's afraid of negative numbers? He'll stop at nothing to avoid them.",
        "Why don't skeletons fight each other? They don't have the guts."
    ]
    speak(random.choice(jokes))


def play_rock_paper_scissors():
    choices = ["rock", "paper", "scissors"]
    speak("Let's play Rock-Paper-Scissors! Say your choice.")
    user_choice = listen().lower()
    if user_choice not in choices:
        speak("Invalid choice. Please say rock, paper, or scissors.")
        return
    ai_choice = random.choice(choices)
    speak(f"I chose {ai_choice}.")
    if user_choice == ai_choice:
        speak("It's a tie!")
    elif (user_choice == "rock" and ai_choice == "scissors") or \
         (user_choice == "paper" and ai_choice == "rock") or \
         (user_choice == "scissors" and ai_choice == "paper"):
        speak("You win!")
    else:
        speak("I win!")

def generate_poetry():
    lines = [
        "Roses are red, violets are blue,",
        "The stars shine bright, just like you.",
        "In the quiet of the night,",
        "Dreams take flight, out of sight.",
        "Whispers of the wind, so soft and low,",
        "Tell tales of places we long to go."
    ]
    poem = "\n".join(random.sample(lines, 4))
    speak("Here's a poem for you:")
    speak(poem)

def generate_meme(top_text, bottom_text):
    speak("Meme generation is not available in offline mode.")

def ask_trivia():
    questions = [
        {
            "question": "What is the capital of France?",
            "answer": "Paris",
            "options": ["London", "Berlin", "Madrid", "Paris"]
        },
        {
            "question": "How many continents are there?",
            "answer": "7",
            "options": ["5", "6", "7", "8"]
        }
    ]
    q = random.choice(questions)
    speak(f"Here's a trivia question: {q['question']}")
    speak(f"Options: {', '.join(q['options'])}")
    user_answer = listen().lower()
    if user_answer == q['answer'].lower():
        speak("Correct! Well done!")
    else:
        speak(f"Sorry, the correct answer is {q['answer']}.")


def calculator():
    """Enhanced scientific calculator with offline capabilities"""
    speak("What would you like to calculate? You can ask about:")
    speak("- Equations (solve 2x + 5 = 15)")
    speak("- Unit conversions (convert 5 feet to meters)")
    speak("- Scientific constants (speed of light)")
    speak("- Statistics (mean of 5,10,15)")
    speak("- Plotting (plot sin(x))")
    
    query = listen()
    if not query:
        return
        
    try:
        # Equation solving
        if "solve" in query:
            eq = query.replace("solve", "").strip()
            x = sp.symbols('x')
            solution = sp.solve(eq, x)
            speak(f"The solution is {solution}")
            
        # Unit conversions
        elif "convert" in query:
            parts = query.replace("convert", "").strip().split(" to ")
            if len(parts) == 2:
                value, from_unit = parts[0].split()
                to_unit = parts[1]
                converted = sp.convert_to(float(value)*sp.Unit(from_unit), sp.Unit(to_unit))
                speak(f"{value} {from_unit} is {converted} {to_unit}")
                
        # Scientific constants
        elif any(word in query for word in ["constant", "speed of light", "gravitational"]):
            constants = {
                "speed of light": "299,792,458 m/s",
                "gravitational constant": "6.67430 × 10⁻¹¹ m³⋅kg⁻¹⋅s⁻²",
                "plank constant": "6.62607015 × 10⁻³⁴ J⋅Hz⁻¹"
            }
            for name, value in constants.items():
                if name in query:
                    speak(f"{name} is {value}")
                    return
                    
        # Statistics
        elif "mean of" in query or "average of" in query:
            nums = [float(n) for n in re.findall(r"\d+", query)]
            if nums:
                speak(f"The mean is {sum(nums)/len(nums):.2f}")
                
        # Plotting
        elif "plot" in query:
            expr = query.replace("plot", "").strip()
            x = sp.symbols('x')
            plot = sp.plot(sp.sympify(expr), (x, -10, 10), show=False)
            plot.save('plot.png')
            speak("I've saved the plot as plot.png")
            
        # Fallback to original calculator
        else:
            # Date calculations
            if "days between" in query or "days from" in query:
                dates = re.findall(r"(\d{1,2}/\d{1,2}/\d{2,4})", query)
                if len(dates) == 2:
                    date1 = datetime.strptime(dates[0], "%m/%d/%Y")
                    date2 = datetime.strptime(dates[1], "%m/%d/%Y")
                    delta = abs((date2 - date1).days)
                    speak(f"There are {delta} days between {dates[0]} and {dates[1]}")
                    return

            # Standard math evaluation
            result = sp.sympify(query).evalf()
            speak(f"The result is {result}")
            
    except Exception as e:
        speak(f"Sorry, I couldn't perform that calculation. Error: {e}")

def wireless_earphone_button_pressed():
    """Check if wireless earphone button is pressed with proper implementation"""
    try:
        # This would interface with actual hardware/Bluetooth in a real implementation
        # For now simulate random button presses for testing
        return random.random() < 0.1  # 10% chance of button press
    except Exception as e:
        print(f"Error checking earphone button: {e}")
        return False


def record_audio_by_button():
    speak("Please press the Play/Pause button on your wireless earphone to start recording and press it again to stop recording.")
    while True:
        if wireless_earphone_button_pressed():
            speak("Recording started. Please press the Play/Pause button again to stop recording.")
            # Start recording audio
            CHUNK = 1024
            FORMAT = pyaudio.paInt16
            CHANNELS = 2
            RATE = 44100
            RECORD_SECONDS = 60
            WAVE_OUTPUT_FILENAME = "output.wav"
            p = pyaudio.PyAudio()
            stream = p.open(format=FORMAT,
                            channels=CHANNELS,
                            rate=RATE,
                            input=True,
                            frames_per_buffer=CHUNK)
            frames = []
            for i in range(0, int(RATE / CHUNK * RECORD_SECONDS)):
                data = stream.read(CHUNK)
                frames.append(data)
                if wireless_earphone_button_pressed():
                    break
            stream.stop_stream()
            stream.close()
            p.terminate()
            wf = wave.open(WAVE_OUTPUT_FILENAME, 'wb')
            wf.setnchannels(CHANNELS)
            wf.setsampwidth(p.get_sample_size(FORMAT))
            wf.setframerate(RATE)
            wf.writeframes(b''.join(frames))
            wf.close()
            speak("Recording stopped and saved as output.wav")
            break

def get_formula_one_updates():
    speak("Formula One updates are not available in offline mode.")


# Cache for search results
search_cache = {}
CACHE_FILE = "search_cache.json"

def load_search_cache():
    """Load search results cache from file"""
    try:
        with open(CACHE_FILE, 'r') as f:
            return json.load(f)
    except (FileNotFoundError, json.JSONDecodeError):
        return {}

def save_search_cache():
    """Save search results cache to file"""
    with open(CACHE_FILE, 'w') as f:
        json.dump(search_cache, f)

def clear_search_cache():
    """Clear the search results cache"""
    global search_cache
    search_cache = {}
    save_search_cache()
    return "Search cache cleared"

from concurrent.futures import ThreadPoolExecutor
import hashlib

def process_file(file_path, file_ext, query, results):
    """Process a single file during offline content search"""
    try:
        if file_ext == '.pdf':
            try:
                with open(file_path, 'rb') as f:
                    reader = PyPDF2.PdfReader(f)
                    for i, page in enumerate(reader.pages):
                        text = page.extract_text()
                        if text and query.lower() in text.lower():
                            results['PDF'].append({
                                'file': file_path,
                                'page': i+1,
                                'text': text[:500] + '...'
                            })
                            if len(results['PDF']) >= 3:
                                break
            except Exception as e:
                print(f"Error reading PDF {file_path}: {e}")
        elif file_ext in ('.png', '.jpg', '.jpeg'):
            try:
                img = Image.open(file_path)
                text = pytesseract.image_to_string(img)
                if query.lower() in text.lower():
                    results['Image'].append({
                        'file': file_path,
                        'text': text[:500] + '...'
                    })
                    if len(results['Image']) >= 3:
                        return
            except Exception as e:
                print(f"Error reading image {file_path}: {e}")
        elif file_ext in ['.txt', '.md', '.csv', '.json', '.xml']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                    if query.lower() in text.lower():
                        results['Text'].append({
                            'file': file_path,
                            'text': text[:500] + '...'
                        })
            except Exception as e:
                print(f"Error reading text file {file_path}: {e}")
        elif file_ext in ['.docx', '.xlsx', '.pptx']:
            try:
                text = ''
                if file_ext == '.docx':
                    try:
                        import docx
                        doc = docx.Document(file_path)
                        text = '\n'.join([para.text for para in doc.paragraphs])
                    except ImportError:
                        speak("python-docx module not installed - skipping Word documents")
                        return
                elif file_ext == '.xlsx':
                    try:
                        import openpyxl
                        wb = openpyxl.load_workbook(file_path)
                        text = ''
                        for sheet in wb.sheetnames:
                            text += f"\nSheet: {sheet}\n"
                            for row in wb[sheet].iter_rows(values_only=True):
                                text += ' '.join(str(cell) for cell in row if cell) + '\n'
                    except ImportError:
                        speak("openpyxl module not installed - skipping Excel documents")
                        return
                elif file_ext == '.pptx':
                    try:
                        from pptx import Presentation
                        prs = Presentation(file_path)
                        text = ''
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text"):
                                    text += shape.text + '\n'
                    except ImportError:
                        speak("python-pptx module not installed - skipping PowerPoint documents")
                        return
                if text and query.lower() in text.lower():
                    results['Office'].append({
                        'file': file_path,
                        'text': text[:500] + '...'
                    })
            except Exception as e:
                print(f"Error reading office file {file_path}: {e}")
        elif file_ext in ['.py', '.js', '.java', '.c', '.cpp', '.html', '.css']:
            try:
                with open(file_path, 'r', encoding='utf-8') as f:
                    text = f.read()
                    if query.lower() in text.lower():
                        results['Code'].append({
                            'file': file_path,
                            'text': text[:500] + '...'
                        })
            except Exception as e:
                print(f"Error reading code file {file_path}: {e}")
        elif file_ext in ['.epub', '.mobi']:
            try:
                if file_ext == '.epub':
                    import epub
                    book = epub.open_epub(file_path)
                    text = ''
                    for item in book.get_items():
                        if item.get_type() == epub.ITEM_DOCUMENT:
                            text += item.get_content().decode('utf-8') + '\n'
                    book.close()
                if query.lower() in text.lower():
                    results['eBooks'].append({
                        'file': file_path,
                        'text': text[:500] + '...'
                    })
            except Exception as e:
                print(f"Error reading eBook {file_path}: {e}")
        elif file_ext in ['.zip', '.rar', '.7z']:
            try:
                import zipfile
                with zipfile.ZipFile(file_path) as z:
                    for name in z.namelist():
                        if name.lower().endswith(('.txt', '.md', '.csv')):
                            with z.open(name) as f:
                                text = f.read().decode('utf-8')
                                if query.lower() in text.lower():
                                    results['Archives'].append({
                                        'file': f"{file_path}/{name}",
                                        'text': text[:500] + '...'
                                    })
                                    break
            except Exception as e:
                print(f"Error reading archive {file_path}: {e}")
    except Exception as e:
        print(f"Error processing file {file_path}: {e}")

    # Filter by requested file types if specified
    if file_types:
        file_types = [ft.lower() for ft in file_types]
        supported_extensions = {k:v for k,v in supported_extensions.items() 
                              if k.lower() in file_types}
    
    # Initialize results with all supported categories
    results = {category: [] for category in supported_extensions}
    results.update({
        'PDF': [],
        'Image': [],
        'Video': []
    })
    
    total_files = 0
    last_progress_time = time.time()
    
    # Search across all drives with enhanced feedback
    search_start = time.time()
    for drive_idx, drive in enumerate(drives):
        drive_start = time.time()
        speak(f"Searching drive {drive} ({drive_idx+1}/{len(drives)})...")
        
        for root, dirs, files in os.walk(drive):
            # Skip excluded directories
            if any(exclude_dir.lower() in root.lower() for exclude_dir in exclude_dirs):
                dirs[:] = []
                continue
                
            # Progress feedback
            if time.time() - last_progress_time > 30:
                elapsed = int(time.time() - search_start)
                speak(f"Searching... {elapsed//60}m {elapsed%60}s elapsed. Found {total_files} files with {sum(len(r) for r in results.values())} matches")
                last_progress_time = time.time()
            # Process files in parallel
            with ThreadPoolExecutor(max_workers=4) as executor:
                futures = []
                for file in files:
                    total_files += 1
                    file_path = os.path.join(root, file)
                    file_ext = os.path.splitext(file)[1].lower()
                    futures.append(executor.submit(process_file, file_path, file_ext, query, results))
                
                # Wait for all futures to complete
                for future in futures:
                    try:
                        future.result()
                    except Exception as e:
                        print(f"Error processing file: {e}")
                
                # PDF files (now under 'PDF' category)
                if file_ext == '.pdf' and 'PDF' in results:
                    try:
                        with open(file_path, 'rb') as f:
                            reader = PyPDF2.PdfReader(f)
                            for i, page in enumerate(reader.pages):
                                text = page.extract_text()
                                if text and query.lower() in text.lower():
                                    results['pdf'].append({
                                        'file': file_path,
                                        'page': i+1,
                                        'text': text[:500] + '...'
                                    })
                                    if len(results['pdf']) >= 3:
                                        break
                    except Exception as e:
                        print(f"Error reading PDF {file_path}: {e}")

                # Image files (now under 'Image' category)
                elif file_ext in ('.png', '.jpg', '.jpeg') and 'Image' in results:
                    try:
                        img = Image.open(file_path)
                        text = pytesseract.image_to_string(img)
                        if query.lower() in text.lower():
                            results['image'].append({
                                'file': file_path,
                                'text': text[:500] + '...'
                            })
                            if len(results['image']) >= 3:
                                break
                    except Exception as e:
                        print(f"Error reading image {file_path}: {e}")

                # Video files (now under 'Video' category)  
                # Text files (now under 'Text' category)
                elif 'Text' in results and file_ext in supported_extensions['Text']:
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            text = f.read()
                            if query.lower() in text.lower():
                                results['text'].append({
                                    'file': file_path,
                                    'text': text[:500] + '...'
                                })
                    except Exception as e:
                        print(f"Error reading text file {file_path}: {e}")
                        
                # Office documents (now under 'Office' category)
                elif 'Office' in results and file_ext in supported_extensions['Office']:
                    try:
                        text = ''
                        if file_ext == '.docx':
                            try:
                                import docx
                                doc = docx.Document(file_path)
                                text = '\n'.join([para.text for para in doc.paragraphs])
                            except ImportError:
                                speak("python-docx module not installed - skipping Word documents")
                                supported_extensions['Office'].remove('.docx')
                                continue
                                
                        elif file_ext == '.xlsx':
                            try:
                                import openpyxl
                                wb = openpyxl.load_workbook(file_path)
                                text = ''
                                for sheet in wb.sheetnames:
                                    text += f"\nSheet: {sheet}\n"
                                    for row in wb[sheet].iter_rows(values_only=True):
                                        text += ' '.join(str(cell) for cell in row if cell) + '\n'
                            except ImportError:
                                speak("openpyxl module not installed - skipping Excel documents")
                                supported_extensions['Office'].remove('.xlsx')
                                continue
                                
                        elif file_ext == '.pptx':
                            try:
                                from pptx import Presentation
                                prs = Presentation(file_path)
                                text = ''
                                for slide in prs.slides:
                                    for shape in slide.shapes:
                                        if hasattr(shape, "text"):
                                            text += shape.text + '\n'
                            except ImportError:
                                speak("python-pptx module not installed - skipping PowerPoint documents")
                                supported_extensions['Office'].remove('.pptx')
                                continue
                                
                        if text and query.lower() in text.lower():
                            results['office'].append({
                                'file': file_path,
                                'text': text[:500] + '...'
                            })
                    except Exception as e:
                        print(f"Error reading office file {file_path}: {e}")
                        
                # Code files (now under 'Code' category)
                elif 'Code' in results and file_ext in supported_extensions['Code']:
                    try:
                        with open(file_path, 'r', encoding='utf-8') as f:
                            text = f.read()
                            if query.lower() in text.lower():
                                results['code'].append({
                                    'file': file_path,
                                    'text': text[:500] + '...'
                                })
                    except Exception as e:
                        print(f"Error reading code file {file_path}: {e}")

                # eBook files
                elif 'eBooks' in results and file_ext in supported_extensions['eBooks']:
                    try:
                        if file_ext == '.epub':
                            import epub
                            book = epub.open_epub(file_path)
                            text = ''
                            for item in book.get_items():
                                if item.get_type() == epub.ITEM_DOCUMENT:
                                    text += item.get_content().decode('utf-8') + '\n'
                            book.close()
                            
                        if query.lower() in text.lower():
                            results['eBooks'].append({
                                'file': file_path,
                                'text': text[:500] + '...'
                            })
                    except Exception as e:
                        print(f"Error processing file: {e}")

    # Format results for output
    response = []
    found_results = False

    for category in sorted(results.keys()):
        if results[category]:
            found_results = True
            category_results = []
            for result in results[category][:5]:  # Limit to 5 results per category
                if 'page' in result:
                    entry = f"- {os.path.basename(result['file'])} (page {result['page']}):\n  {result['text']}"
                else:
                    entry = f"- {os.path.basename(result['file'])}:\n  {result['text']}"
                category_results.append(entry)

            if category_results:
                response.append(f"=== {category.upper()} ===")
                response.extend(category_results)

    if not found_results:
        final_response = "No results found in your files."
    else:
        final_response = "Search Results:\n\n" + "\n\n".join(response)


# Note-taking feature
def take_notes():
    speak("Please say what you want to note.")
    note = listen()
    with open("notes.txt", "a") as f:
        f.write(note + "\n")
    speak("Note taken.")

# Text-to-speech feature
def text_to_speech():
    speak("Please say what you want to read.")
    text = listen()
    engine.say(text)
    engine.runAndWait()

# Translation feature
def translate(text):
    # Simple translation example - would need proper translation API in production
    translations = {
        "hello": "hola",
        "goodbye": "adiós",
        "thank you": "gracias"
    }
    return translations.get(text.lower(), "Translation not available")

def translate_command():
    speak("Please say what you want to translate.")
    text = listen()
    translation = translate(text)
    speak(translation)


# Math problem-solving feature
def solve_math_problem(problem):
    try:
        return str(eval(problem))
    except:
        return "Could not solve the problem"

def solve_math_command():
    speak("Please say what math problem you want to solve.")
    problem = listen()
    solution = solve_math_problem(problem)
    speak(solution)


# Flashcard feature
def create_flashcard():
    speak("Please say what you want to put on the front of the flashcard.")
    front = listen()
    speak("Please say what you want to put on the back of the flashcard.")
    back = listen()
    with open("flashcards.txt", "a") as f:
        f.write(front + ":" + back + "\n")
    speak("Flashcard created.")

# Audio recording feature
def record_audio():
    speak("Please say what you want to record.")
    recording = listen()
    with open("recording.wav", "wb") as f:
        f.write(recording)
    speak("Recording saved.")

# Voice-to-text feature
def voice_to_text():
    speak("Please say what you want to write.")
    text = listen()
    with open("text.txt", "w") as f:
        f.write(text)
    speak("Text written.")

# Dictionary feature
def look_up_word(word):
    # Basic dictionary implementation
    dictionary = {
        "hello": "a greeting",
        "world": "the earth with all its countries and peoples",
        "python": "a high-level programming language"
    }
    return dictionary.get(word.lower(), "Definition not found")

def look_up_word_command():
    speak("Please say what word you want to look up.")
    word = listen()
    definition = look_up_word(word)
    speak(definition)


# Thesaurus feature
def look_up_synonyms(word):
    # Basic thesaurus implementation
    thesaurus = {
        "happy": ["joyful", "cheerful", "content"],
        "sad": ["unhappy", "depressed", "melancholy"]
    }
    return thesaurus.get(word.lower(), ["No synonyms found"])

def look_up_synonyms_command():
    speak("Please say what word you want synonyms for.")
    word = listen()
    synonyms = look_up_synonyms(word)
    speak(f"Synonyms for {word}: {', '.join(synonyms)}")


# Study guide feature
def create_study_guide():
    speak("Please say what you want to put in the study guide.")
    guide = listen()
    with open("study_guide.txt", "w") as f:
        f.write(guide)
    speak("Study guide created.")

def remember_fact(fact, category="general"):
    """Store a personal fact in the knowledge base"""
    if category not in knowledge_base:
        knowledge_base[category] = []
    knowledge_base[category].append(fact)
    with open(knowledge_file, 'w') as f:
        json.dump(knowledge_base, f)
    return f"I'll remember that: {fact}"

def recall_fact(topic):
    """Retrieve facts about a topic from the knowledge base"""
    results = []
    for category, facts in knowledge_base.items():
        if topic.lower() in category.lower():
            results.extend(facts)
        for fact in facts:
            if topic.lower() in fact.lower():
                results.append(fact)
    return results if results else "I don't know anything about that."

def buzz(command):
    """Process and enhance the recognized command for better interpretation"""
    if not command:
        return None
    command = command.lower().strip()
    
    # Example enhancements:
    # - Replace synonyms or common phrases
    # - Normalize commands
    # - Remove filler words
    # - Expand contractions
    
    # Replace some common synonyms or phrases
    replacements = {
        "turn on": "open",
        "turn off": "close",
        "shut down": "shutdown",
        "restart computer": "restart",
        "lock screen": "lock",
        "sleep mode": "sleep",
        "hibernate mode": "hibernate",
        "add reminder": "set reminder",
        "create reminder": "set reminder",
        "set alarm clock": "set alarm",
        "play music": "open music",
        "stop music": "close music",
        "launch application": "open",
        "quit application": "close",
        "exit application": "close",
        "kill process": "kill",
        "end process": "kill",
        "show me": "show",
        "list all": "list",
        "what time is it": "time",
        "what date is it": "date",
        "calculate this": "calculate",
        "solve this": "calculate",
        "tell me a joke": "joke",
        "tell me something funny": "joke",
        "wake me up": "wake up",
        "go to sleep": "sleep",
        "pause now": "pause",
        "resume now": "resume"
    }
    
    for phrase, replacement in replacements.items():
        if phrase in command:
            command = command.replace(phrase, replacement)
    
    # Remove filler words
    filler_words = ["please", "could you", "would you", "can you", "hey", "ok", "okay"]
    for filler in filler_words:
        command = command.replace(filler, "")
    
    # Remove extra spaces
    command = " ".join(command.split())
    
    return command

def interpret_command(command):
    """Interpret user command with enhanced system control"""
    if not command:
        return None

    try:
        # Pass command through buzz processing first
        command = buzz(command)
        
        # Enhanced command mapping with more comprehensive control
        command_map = {
            'open browser': ('open', 'chrome'),
            'open notepad': ('open', 'notepad'),
            'open calculator': ('open', 'calculator'),
            'open word': ('open', 'winword'),
            'open excel': ('open', 'excel'),
            'open powerpoint': ('open', 'powerpnt'),
            'open outlook': ('open', 'outlook'),
            'open audacity': ('open', 'audacity'),
            'open chrome': ('open', 'chrome'),
            'open firefox': ('open', 'firefox'),
            'open edge': ('open', 'msedge'),
            'open photoshop': ('open', 'photoshop'),
            'open illustrator': ('open', 'illustrator'),
            'open vscode': ('open', 'code'),
            'open sublime': ('open', 'sublime_text'),
            'open pycharm': ('open', 'pycharm'),
            'open eclipse': ('open', 'eclipse'),
            'open intellij': ('open', 'idea'),
            'open steam': ('open', 'steam'),
            'open discord': ('open', 'discord'),
            'open spotify': ('open', 'spotify'),
            'open vlc': ('open', 'vlc'),
            'open adobe reader': ('open', 'acrord32'),
            'describe yourself': 'describe_personality',
            'system information': 'system_info',
            'cpu usage': 'cpu_usage',
            'memory usage': 'memory_usage',
            'disk usage': 'disk_usage',
            'network status': 'network_status',
            'battery level': 'battery_level',
            'increase volume': 'volume_up',
            'decrease volume': 'volume_down',
            'mute volume': 'volume_mute',
            'unmute volume': 'volume_unmute',
            'take screenshot': 'screenshot',
            'capture screen': 'screenshot',
            'record screen': 'screen_record',
            'start recording': 'screen_record',
            'stop recording': 'stop_record',
            'open file manager': 'open_explorer',
            'open explorer': 'open_explorer',
            'open downloads': 'open_downloads',
            'open documents': 'open_documents',
            'open pictures': 'open_pictures',
            'open music': 'open_music',
            'open videos': 'open_videos',
            'search files': 'search_files',
            'find file': 'search_files',
            'send email': 'send_email',
            'compose email': 'send_email',
            'create folder': 'create_folder',
            'make directory': 'create_folder',
            'delete folder': 'delete_folder',
            'remove folder': 'delete_folder',
            'copy file': 'copy_file',
            'paste file': 'paste_file',
            'cut file': 'cut_file',
            'rename file': 'rename_file',
            'move file': 'move_file',
            'compress file': 'compress_file',
            'extract file': 'extract_file',
            'create shortcut': 'create_shortcut',
            'make shortcut': 'create_shortcut',
            'open control panel': 'open_control_panel',
            'open settings': 'open_settings',
            'open task manager': 'open_task_manager',
            'open command prompt': 'open_cmd',
            'open terminal': 'open_cmd',
            'open powershell': 'open_powershell',
            'open registry': 'open_registry',
            'open services': 'open_services',
            'open event viewer': 'open_event_viewer',
            'open device manager': 'open_device_manager',
            'open network settings': 'open_network_settings',
            'open firewall': 'open_firewall',
            'open antivirus': 'open_antivirus',
            'check antivirus': 'check_antivirus',
            'update system': 'update_system',
            'check updates': 'update_system',
            'install updates': 'update_system',
            'create backup': 'create_backup',
            'backup files': 'create_backup',
            'restore backup': 'restore_backup',
            'restore files': 'restore_backup',
            'schedule task': 'schedule_task',
            'create task': 'schedule_task',
            'cancel task': 'cancel_task',
            'delete task': 'cancel_task',
            'open task scheduler': 'open_task_scheduler',
            'open performance monitor': 'open_perfmon',
            'open resource monitor': 'open_resmon',
            'open disk cleanup': 'open_cleanup',
            'open disk defragmenter': 'open_defrag',
            'defragment disk': 'open_defrag',
            'open system restore': 'open_restore',
            'restore system': 'open_restore',
            'open backup and restore': 'open_backup_restore',
            'open file history': 'open_file_history',
            'open storage sense': 'open_storage_sense',
            'open storage settings': 'open_storage_settings',
            'open power options': 'open_power_options',
            'open display settings': 'open_display_settings',
            'open sound settings': 'open_sound_settings',
            'open network and sharing': 'open_network_sharing',
            'open bluetooth settings': 'open_bluetooth',
            'open printer settings': 'open_printers',
            'open mouse settings': 'open_mouse_settings',
            'open keyboard settings': 'open_keyboard_settings',
            'open date and time': 'open_date_time',
            'open user accounts': 'open_user_accounts',
            'open system properties': 'open_system_properties',
            'open advanced system settings': 'open_advanced_system',
            'open environment variables': 'open_env_vars',
            'open startup programs': 'open_startup',
            'open startup folder': 'open_startup_folder',
            'open fonts folder': 'open_fonts',
            'open system32': 'open_system32',
            'open temp folder': 'open_temp',
            'open recycle bin': 'open_recycle_bin',
            'empty recycle bin': 'empty_recycle_bin',
            'open clipboard': 'open_clipboard',
            'clear clipboard': 'clear_clipboard',
            'open magnifier': 'open_magnifier',
            'open narrator': 'open_narrator',
            'open onscreen keyboard': 'open_osk',
            'open sticky keys': 'open_sticky_keys',
            'open filter keys': 'open_filter_keys',
            'open high contrast': 'open_high_contrast',
            'open ease of access': 'open_ease_of_access',
            'open accessibility': 'open_ease_of_access',
        }

        # Return the command map for external use
        return command_map
    except Exception as e:
        print(f"Error in interpret_command: {e}")
        return None

def control_smart_home(device, action):
    """Control smart home devices with enhanced capabilities"""
    # Lights control
    if device in SMART_HOME_DEVICES.get('lights', {}):
        ip = SMART_HOME_DEVICES['lights'][device]
        try:
            requests.post(f"http://{ip}/control", 
                        json={'action': action},
                        timeout=3)
            speak(f"Turning {action} {device} lights")
            return True
        except Exception as e:
            print(f"Smart home error: {e}")
            speak(f"Failed to control {device} lights")
            return False
            
    # Thermostat control        
    elif device == 'thermostat' and (action.isdigit() or action in ['heat', 'cool', 'off']):
        ip = SMART_HOME_DEVICES['thermostat']
        try:
            if action.isdigit():
                requests.post(f"http://{ip}/set", 
                            json={'temp': int(action)},
                            timeout=3)
                speak(f"Setting thermostat to {action} degrees")
            else:
                requests.post(f"http://{ip}/mode", 
                            json={'mode': action},
                            timeout=3)
                speak(f"Setting thermostat mode to {action}")
            return True
        except Exception as e:
            print(f"Thermostat error: {e}")
            speak("Failed to adjust thermostat")
            return False
            
    # Security system control
    elif device == 'security' and action in ['arm', 'disarm']:
        if 'security' in SMART_HOME_DEVICES:
            ip = SMART_HOME_DEVICES['security']
            try:
                requests.post(f"http://{ip}/{action}",
                            timeout=3)
                speak(f"Security system {action}ed")
                return True
            except Exception as e:
                print(f"Security system error: {e}")
                speak(f"Failed to {action} security system")
                return False
                
    return False

def sync_with_cloud():
    """Sync data with cloud service"""
    if not CLOUD_SYNC_ENABLED:
        return False
    
    data = {
        'todos': todo_list,
        'reminders': reminders,
        'context': context_memory[-10:] if context_memory else []
    }
    
    try:
        response = requests.post(
            f"{CLOUD_CONFIG['cloud_url']}/sync",
            json=data,
            headers={'Authorization': f"Bearer {CLOUD_CONFIG['api_key']}"}
        )
        return response.status_code == 200
    except:
        return False

import spacy
import threading
import queue
import logging
import hashlib
import json
import time

# Initialize spaCy NLP model
try:
    nlp = spacy.load("en_core_web_sm")
except OSError:
    print("spaCy model not found. Installing...")
    import subprocess
    subprocess.run(["python", "-m", "spacy", "download", "en_core_web_sm"])
    nlp = spacy.load("en_core_web_sm")

# Custom intents for better command recognition
INTENTS = {
    "greeting": ["hello", "hi", "hey", "good morning", "good afternoon", "good evening"],
    "farewell": ["bye", "goodbye", "see you", "farewell"],
    "time": ["what time", "current time", "time now"],
    "date": ["what date", "current date", "today's date"],
    "weather": ["weather", "forecast", "temperature"],
    "joke": ["tell joke", "joke", "funny"],
    "calculator": ["calculate", "math", "compute"],
    "open_app": ["open", "launch", "start"],
    "close_app": ["close", "quit", "exit"],
    "alarm": ["set alarm", "alarm"],
    "reminder": ["remind", "reminder"],
    "task": ["add task", "task", "todo"],
    "help": ["help", "assist", "what can you do"],
    "standby": ["standby", "sleep", "rest"],
    "wake": ["wake up", "activate"],
    "personality": ["describe yourself", "who are you"],
    "volume": ["volume up", "volume down", "mute", "unmute"],
    "system": ["system info", "cpu", "memory", "disk"],
    "search": ["search", "find"],
    "note": ["take note", "note"],
    "translate": ["translate"],
    "solve": ["solve", "math problem"]
}

def classify_intent(text):
    """Classify user input into predefined intents using keyword matching and NLP"""
    doc = nlp(text.lower())
    scores = {}

    # Keyword-based intent classification
    for intent, keywords in INTENTS.items():
        score = 0
        for keyword in keywords:
            if keyword in text.lower():
                score += 1
        if score > 0:
            scores[intent] = score

    # Use NLP entities for additional context
    entities = [(ent.text, ent.label_) for ent in doc.ents]
    if entities:
        for entity_text, label in entities:
            if label == "TIME":
                scores["time"] = scores.get("time", 0) + 1
            elif label == "DATE":
                scores["date"] = scores.get("date", 0) + 1

    # Return intent with highest score
    if scores:
        return max(scores, key=scores.get)
    return "unknown"

# Dialog Manager class
import json

class DialogManager:
    def __init__(self, context_file="dialog_context.json"):
        self.context_file = context_file
        self.context = self.load_context()
        self.max_context_length = 10

    def update_context(self, user_input, system_response):
        self.context.append({"user": user_input, "system": system_response})
        if len(self.context) > self.max_context_length:
            self.context.pop(0)
        self.save_context()

    def get_context(self):
        return self.context

    def clear_context(self):
        self.context = []
        self.save_context()

    def save_context(self):
        try:
            with open(self.context_file, "w") as f:
                json.dump(self.context, f)
        except Exception as e:
            print(f"Error saving dialog context: {e}")

    def load_context(self):
        try:
            with open(self.context_file, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            return []

# Personalization module
class Personalization:
    def __init__(self, filename="personalization.json"):
        self.filename = filename
        self.data = self.load_data()

    def load_data(self):
        try:
            with open(self.filename, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            return {}

    def save_data(self):
        with open(self.filename, "w") as f:
            json.dump(self.data, f)

    def remember(self, key, value):
        self.data[key] = value
        self.save_data()

    def recall(self, key):
        return self.data.get(key, None)

# Multimodal input handler (stub)
class MultimodalInputHandler:
    def __init__(self):
        pass

    def process_voice(self, voice_input):
        # Process voice input (already handled)
        return voice_input

    def process_text(self, text_input):
        # Process text input
        return text_input

    def process_image(self, image_input):
        # Process image input (stub)
        return "Image processed"

# Feedback mechanism
class FeedbackManager:
    def __init__(self, filename="feedback.json"):
        self.filename = filename
        self.feedback_list = self.load_feedback()

    def load_feedback(self):
        try:
            with open(self.filename, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            return []

    def save_feedback(self):
        with open(self.filename, "w") as f:
            json.dump(self.feedback_list, f)

    def add_feedback(self, feedback):
        self.feedback_list.append({"timestamp": time.time(), "feedback": feedback})
        self.save_feedback()

# Enhanced output mechanism
class OutputManager:
    def __init__(self):
        self.log_file = "interaction.log"
        logging.basicConfig(filename=self.log_file, level=logging.INFO,
                            format='%(asctime)s - %(levelname)s - %(message)s')

    def output(self, text, speak_flag=True):
        if speak_flag:
            speak(text)
        print(text)
        logging.info(text)

# Analytics and logs
class AnalyticsManager:
    def __init__(self):
        self.command_count = 0
        self.start_time = time.time()

    def log_command(self, command):
        self.command_count += 1
        logging.info(f"Command received: {command}")

    def get_usage_stats(self):
        uptime = time.time() - self.start_time
        return {"commands_processed": self.command_count, "uptime_seconds": uptime}

# Multitasking feature using threading
class TaskManager:
    def __init__(self):
        self.task_queue = queue.Queue()
        self.worker_thread = threading.Thread(target=self.worker)
        self.worker_thread.daemon = True
        self.worker_thread.start()

    def worker(self):
        while True:
            task_func, args = self.task_queue.get()
            try:
                task_func(*args)
            except Exception as e:
                logging.error(f"Error in task: {e}")
            self.task_queue.task_done()
multimodal_handler = MultimodalInputHandler()
feedback_manager = FeedbackManager()
output_manager = OutputManager()
analytics_manager = AnalyticsManager()
task_manager = TaskManager()
security_manager = SecurityManager()

# Personality Manager for AI personality refinement
class PersonalityManager:
    def __init__(self, filename="personality.json"):
        self.filename = filename
        self.personality_traits = self.load_personality()

    def load_personality(self):
        try:
            with open(self.filename, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            # Default personality traits
            return {
                "name": "Vicky",
                "tone": "friendly",
                "humor": True,
                "formality": "casual"
            }

    def save_personality(self):
        with open(self.filename, "w") as f:
            json.dump(self.personality_traits, f)

    def update_trait(self, trait, value):
        self.personality_traits[trait] = value
        self.save_personality()

    def get_trait(self, trait):
        return self.personality_traits.get(trait, None)

    def describe_personality(self):
        desc = f"My name is {self.personality_traits.get('name', 'Assistant')}. "
        desc += f"I speak in a {self.personality_traits.get('tone', 'neutral')} tone. "
        desc += f"My humor is {'enabled' if self.personality_traits.get('humor', False) else 'disabled'}. "
        desc += f"My formality level is {self.personality_traits.get('formality', 'neutral')}."
        return desc

# Customization Hub for user preferences
class CustomizationHub:
    def __init__(self, filename="customization.json"):
        self.filename = filename
        self.preferences = self.load_preferences()

    def load_preferences(self):
        try:
            with open(self.filename, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            # Default preferences
            return {
                "voice_speed": 150,
                "voice_pitch": 0.5,
                "language": "en",
                "theme": "default"
            }

    def save_preferences(self):
        with open(self.filename, "w") as f:
            json.dump(self.preferences, f)

    def update_preference(self, key, value):
        self.preferences[key] = value
        self.save_preferences()

    def get_preference(self, key):
        return self.preferences.get(key, None)

# Problem Solver for dynamic problem solving
class ProblemSolver:
    def __init__(self):
        pass

    def solve(self, query, context=None):
        # Placeholder for integration with AI reasoning API or local model
        # For now, simple echo or basic math evaluation
        try:
            # Try to evaluate as math expression
            result = sp.sympify(query).evalf()
            return f"The result is {result}"
        except:
            # Fallback response
            return "I'm working on solving that. This feature will be enhanced soon."

# Advanced User Profiles
class UserProfileManager:
    def __init__(self, filename="user_profiles.json"):
        self.filename = filename
        self.profiles = self.load_profiles()
        self.current_user = None

    def load_profiles(self):
        try:
            with open(self.filename, "r") as f:
                return json.load(f)
        except (FileNotFoundError, json.JSONDecodeError):
            return {}

    def save_profiles(self):
        with open(self.filename, "w") as f:
            json.dump(self.profiles, f)

    def add_user(self, username, data=None):
        if username not in self.profiles:
            self.profiles[username] = data or {}
            self.save_profiles()

    def set_current_user(self, username):
        if username in self.profiles:
            self.current_user = username
            return True
        return False

    def get_current_user_data(self):
        if self.current_user:
            return self.profiles.get(self.current_user, {})
        return {}

    def update_user_data(self, key, value):
        if self.current_user:
            self.profiles[self.current_user][key] = value
            self.save_profiles()

# Conversational Expansions
class ConversationalExpander:
    def __init__(self):
        pass

    def expand(self, input_text, context=None):
        # Placeholder for advanced conversational logic
        # For now, simple echo with context awareness
        if context:
            return f"Continuing our conversation: {input_text}"
        else:
            return input_text

# Real-Time Analytics (extended)
class RealTimeAnalyticsManager(AnalyticsManager):
    def __init__(self):
        super().__init__()
        self.live_data = []

    def log_command(self, command):
        super().log_command(command)
        self.live_data.append({"timestamp": time.time(), "command": command})

    def get_live_data(self):
        return self.live_data[-100:]  # Last 100 commands

# Self-Learning Capabilities
class LearningManager:
    def __init__(self, feedback_manager, analytics_manager):
        self.feedback_manager = feedback_manager
        self.analytics_manager = analytics_manager

    def analyze_feedback(self):
        # Placeholder for learning from feedback
        feedbacks = self.feedback_manager.feedback_list
        # Analyze feedback to improve responses (not implemented)
        return len(feedbacks)

    def adapt_behavior(self):
        # Placeholder for adapting behavior based on analytics
        stats = self.analytics_manager.get_usage_stats()
        # Adapt personality or responses (not implemented)
        return stats

# Device Linking
class DeviceLinkManager:
    def __init__(self):
        self.linked_devices = []

    def link_device(self, device_id):
        if device_id not in self.linked_devices:
            self.linked_devices.append(device_id)
            return True
        return False

    def unlink_device(self, device_id):
        if device_id in self.linked_devices:
            self.linked_devices.remove(device_id)
            return True
        return False

    def get_linked_devices(self):
        return self.linked_devices

# Device Control Manager with multi-device support
class DeviceControlManager:
    def __init__(self):
        self.linked_devices = device_link_manager.get_linked_devices()

    def execute_command(self, device, command, params=None):
        """
        Execute a command on a specified device.
        device: device identifier or type (e.g., 'pc', 'phone', 'speaker', 'tv')
        command: command string (e.g., 'open app', 'volume up', 'play music')
        params: optional parameters dict
        """
        device = device.lower()
        command = command.lower()
        if device == 'pc':
            return self._execute_pc_command(command, params)
        elif device == 'phone':
            return self._execute_phone_command(command, params)
        elif device == 'speaker':
            return self._execute_speaker_command(command, params)
        elif device == 'tv':
            return self._execute_tv_command(command, params)
        else:
            return f"Device '{device}' not supported."

    def _execute_pc_command(self, command, params):
        # Basic PC command execution using subprocess or OS commands
        try:
            import platform
            system = platform.system()
            if system != 'Windows':
                return f"PC command '{command}' not implemented for {system}."

            import os
            import ctypes
            import comtypes
            from ctypes import POINTER, cast
            from comtypes import CLSCTX_ALL
            from pycaw.pycaw import AudioUtilities, IAudioEndpointVolume

            def set_volume(level):
                devices = AudioUtilities.GetSpeakers()
                interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                volume = cast(interface, POINTER(IAudioEndpointVolume))
                volume.SetMasterVolumeLevelScalar(level, None)

            def change_volume(delta):
                devices = AudioUtilities.GetSpeakers()
                interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                volume = cast(interface, POINTER(IAudioEndpointVolume))
                current = volume.GetMasterVolumeLevelScalar()
                new_volume = min(max(0.0, current + delta), 1.0)
                volume.SetMasterVolumeLevelScalar(new_volume, None)

            if command.startswith('open '):
                app_name = command[5:].strip()
                return open_application(app_name)
            elif command in ['shutdown', 'restart', 'lock', 'sleep']:
                # Map to system commands
                import platform
                system = platform.system()
                if system == 'Windows':
                    import os
                    if command == 'shutdown':
                        os.system('shutdown /s /t 1')
                    elif command == 'restart':
                        os.system('shutdown /r /t 1')
                    elif command == 'lock':
                        import ctypes
                        ctypes.windll.user32.LockWorkStation()
                    elif command == 'sleep':
                        os.system('rundll32.exe powrprof.dll,SetSuspendState 0,1,0')
                    return f"PC command '{command}' executed."
                else:
                    return f"PC command '{command}' not implemented for {system}."
            elif command.startswith('close '):
                app_name = command[6:].strip()
                os.system(f'taskkill /IM {app_name}.exe /F')
                return f"Closed {app_name} on PC."
            elif command == 'volume up':
                change_volume(0.1)
                return "Volume increased."
            elif command == 'volume down':
                change_volume(-0.1)
                return "Volume decreased."
            elif command == 'mute':
                devices = AudioUtilities.GetSpeakers()
                interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                volume = cast(interface, POINTER(IAudioEndpointVolume))
                volume.SetMute(1, None)
                return "Volume muted."
            elif command == 'unmute':
                devices = AudioUtilities.GetSpeakers()
                interface = devices.Activate(IAudioEndpointVolume._iid_, CLSCTX_ALL, None)
                volume = cast(interface, POINTER(IAudioEndpointVolume))
                volume.SetMute(0, None)
                return "Volume unmuted."
            else:
                return f"PC command '{command}' not recognized."
        except Exception as e:
            return f"Error executing PC command: {e}"
 
    def _execute_phone_command(self, command, params):
        # Placeholder for phone command execution
        # Could integrate with Android intents or iOS shortcuts via external APIs
        return f"Phone command '{command}' executed (stub)."

    def _execute_speaker_command(self, command, params):
        # Placeholder for smart speaker commands (e.g., volume, play, pause)
        if command in ['volume up', 'volume down', 'mute', 'unmute', 'play', 'pause', 'stop']:
            return f"Speaker command '{command}' executed (stub)."
        else:
            return f"Speaker command '{command}' not recognized."

    def _execute_tv_command(self, command, params):
        # Placeholder for smart TV commands (e.g., power on/off, channel change)
        if command in ['power on', 'power off', 'volume up', 'volume down', 'mute', 'unmute', 'channel up', 'channel down']:
            return f"TV command '{command}' executed (stub)."
        else:
            return f"TV command '{command}' not recognized."


# Initialize new managers
personality_manager = PersonalityManager()
customization_hub = CustomizationHub()
problem_solver = ProblemSolver()
user_profile_manager = UserProfileManager()
conversational_expander = ConversationalExpander()
real_time_analytics_manager = RealTimeAnalyticsManager()
learning_manager = LearningManager(feedback_manager, analytics_manager)
device_link_manager = DeviceLinkManager()
device_control_manager = DeviceControlManager()


def main():
    global is_standby
    global race_mode_active
    global lite_mode


    # Interactive setup wizard on first launch
    first_launch = get_user_preference("first_launch")
    if first_launch is None:
        speak("Hey! I’m Bough. Let’s get you set up, yeah?")
        speak("First, let's test your microphone. Please say something.")
        test_input = listen()
        if test_input:
            speak(f"I heard: {test_input}. Setup complete!")
            save_user_preference("first_launch", "completed")
        else:
            speak("I didn't catch that. Let's try again later.")
            save_user_preference("first_launch", "incomplete")

    speak("Initializing all systems. Online and ready.", "alert")

    # Initialize SimRacingAssistant instance
    sim_racing_assistant = SimRacingAssistant()

    # Initialize race mode state
    race_mode_active = False

    while True:
        try:
            if not is_standby:
                command = listen()
                if command:
                    action = interpret_command(command)

                    # Check for race mode activation/deactivation commands
                    if "activate race mode" in command.lower() or "start race mode" in command.lower():
                        race_mode_active = True
                        continue
                    elif "deactivate race mode" in command.lower() or "stop race mode" in command.lower():
                        race_mode_active = False
                        continue

                    # If race mode is active, handle sim racing commands
                    if race_mode_active and any(keyword in command.lower() for keyword in ["car setup", "tyre", "tire", "gap", "strategy", "engineer", "advisor", "team", "telemetry", "radio", "weather", "opponent", "pit wall", "preset strategy"]):
                        response = sim_racing_assistant.assist(command)
                        speak(response)
                    elif action == 'open':
                        app_name = command.replace("open", "").strip()
                        if app_name:
                            open_application(app_name)
                    elif action == 'set_alarm':
                        alarm_time_str = listen()
                        if alarm_time_str:
                            set_alarm(alarm_time_str)
                    elif action == 'add_task':
                        task = listen()
                        if task:
                            add_task(task)
                    elif action == 'show_tasks':
                        show_tasks()
                    elif action == 'set_reminder':
                        reminder_details = listen()
                        if reminder_details:
                            set_reminder(reminder_details, "12:00 PM")
                    elif action == 'tell_joke':
                        tell_joke()
                    elif action == 'get_time':
                        speak(f"The time is {datetime.now().strftime('%I:%M %p')}")
                    elif action == 'get_date':
                        speak(f"Today is {datetime.now().strftime('%B %d, %Y')}")
                    elif action == 'calculator':
                        calculator()
                    elif action == 'formula_one_updates':
                        get_formula_one_updates()
                    elif action == 'standby':
                        is_standby = True
                    elif action == 'wake_up':
                        is_standby = False
                    elif action == 'help':
                        speak("I can help you with tasks, reminders, jokes, calculations, and more.")
                    elif action == 'describe_personality':
                        desc = personality_manager.describe_personality()
                        speak(desc)
                    elif action == 'open_audacity':
                        open_audacity()
                    elif action is None:
                        fallback_response = respond(command if command else "")
                        speak(fallback_response)
                    else:
                        speak(f"Command {action} is not implemented yet.")
                else:
                    speak("I didn't catch that. Please repeat.")
            else:
                time.sleep(1)
        except KeyboardInterrupt:
            speak("Shutting down. Goodbye!")
            break
        except Exception as e:
            speak(f"An error occurred: {e}")
            time.sleep(1)

def recognize_speech_google():
    """
    Listens for a command and attempts to recognize it using Google Speech Recognition.
    """
    r = sr.Recognizer()
    r.energy_threshold = 4000  # Adjust energy threshold for better performance
    with sr.Microphone() as source:
        print("Say something!")
        # Adjust for ambient noise to improve recognition accuracy
        r.adjust_for_ambient_noise(source)
        try:
            audio = r.listen(source)
            print("Recognizing...")
            # Use Google Speech Recognition
            text = r.recognize_google(audio)
            print(f"You said: \"{text}\"")
            return text
        except sr.UnknownValueError:
            print("Google Speech Recognition could not understand audio.")
            return None
        except sr.RequestError as e:
            print(f"Could not request results from Google Speech Recognition service; {e}")
            return None
        except Exception as e:
            print(f"An unexpected error occurred: {e}")
            return None

if __name__ == "__main__":
    main()

#audio=r.listen
